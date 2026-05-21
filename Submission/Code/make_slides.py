"""
Generate RACE Quiz System — Demo Presentation (.pptx)
Run: python make_slides.py
Output: notebooks/RACE_Quiz_System_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ────────────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x1A, 0x1A, 0x2E)   # near-black navy
MID_BG     = RGBColor(0x16, 0x21, 0x3E)   # dark blue
ACCENT     = RGBColor(0x0F, 0x3A, 0x8C)   # deep blue accent
HIGHLIGHT  = RGBColor(0x00, 0xB4, 0xD8)   # cyan highlight
GREEN      = RGBColor(0x06, 0xD6, 0xA0)   # teal/green
ORANGE     = RGBColor(0xFF, 0xA0, 0x26)   # amber
RED        = RGBColor(0xEF, 0x47, 0x6F)   # coral red
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xD6, 0xF1)
SUBTEXT    = RGBColor(0x8F, 0xA3, 0xC8)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # truly blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color=DARK_BG):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape


def add_textbox(slide, text, l, t, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_para(tf, text, font_size=16, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p


def slide_number(slide, n, total, prs):
    add_textbox(slide, f"{n} / {total}",
                Inches(12.3), Inches(7.1), Inches(1), Inches(0.3),
                font_size=11, color=SUBTEXT, align=PP_ALIGN.RIGHT)


def accent_bar(slide, color=HIGHLIGHT, height=Inches(0.07)):
    add_rect(slide, 0, 0, SLIDE_W, height, color)


def section_header_bar(slide, text, color=ACCENT):
    add_rect(slide, 0, Inches(1.05), SLIDE_W, Inches(0.06), color)
    add_textbox(slide, text,
                Inches(0.45), Inches(0.62), Inches(12), Inches(0.55),
                font_size=28, bold=True, color=HIGHLIGHT, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    # gradient-like left panel
    add_rect(s, 0, 0, Inches(5.2), SLIDE_H, MID_BG)
    add_rect(s, 0, 0, Inches(0.18), SLIDE_H, HIGHLIGHT)

    add_textbox(s, "RACE", Inches(0.4), Inches(1.4), Inches(4.6), Inches(1.1),
                font_size=72, bold=True, color=HIGHLIGHT, align=PP_ALIGN.LEFT)
    add_textbox(s, "Intelligent Reading Comprehension\n& Quiz Generation System",
                Inches(0.4), Inches(2.5), Inches(4.8), Inches(1.4),
                font_size=22, bold=False, color=WHITE, align=PP_ALIGN.LEFT)
    add_textbox(s, "AI Lab Project — Spring 2026\nFAST NUCES BS CS",
                Inches(0.4), Inches(4.2), Inches(4.8), Inches(0.8),
                font_size=15, color=SUBTEXT, align=PP_ALIGN.LEFT)

    # right panel — system overview bullets
    add_textbox(s, "System at a Glance",
                Inches(5.6), Inches(1.3), Inches(7.3), Inches(0.55),
                font_size=20, bold=True, color=HIGHLIGHT, align=PP_ALIGN.LEFT)

    items = [
        "Dataset: RACE — 28,000 passages, 100K+ MCQs",
        "Model A: Answer Verifier + Question Generator",
        "Model B: Distractor Generator + Hint Scorer",
        "Metrics: BLEU-4, ROUGE-1/2/L, METEOR",
        "UI: 4-screen Streamlit application",
        "Approach: Traditional ML (LinearSVC, RF, LR)",
    ]
    txb = slide.shapes.add_textbox if False else \
          s.shapes.add_textbox(Inches(5.6), Inches(2.0), Inches(7.3), Inches(4.5))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = item
        run.font.size  = Pt(17)
        run.font.color.rgb = LIGHT_GREY
    return s


def slide_toc(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    add_textbox(s, "Presentation Outline", Inches(0.45), Inches(0.22), Inches(12), Inches(0.6),
                font_size=30, bold=True, color=WHITE)
    add_rect(s, 0, Inches(1.0), SLIDE_W, Inches(0.05), HIGHLIGHT)

    sections = [
        ("01", "Dataset Overview & Structure",        "p. 3–4"),
        ("02", "Exploratory Data Analysis (EDA)",     "p. 5–7"),
        ("03", "Statistical Analysis",                "p. 8–9"),
        ("04", "Data Preprocessing Pipeline",         "p. 10–11"),
        ("05", "Model A — Verifier & QGen",           "p. 12–16"),
        ("06", "Model B — Distractor & Hints",        "p. 17–20"),
        ("07", "Evaluation Results & Interpretation", "p. 21–23"),
        ("08", "System Architecture & UI",            "p. 24–25"),
    ]
    cols = [sections[:4], sections[4:]]
    for ci, col in enumerate(cols):
        x = Inches(0.5 + ci * 6.4)
        for ri, (num, title, pg) in enumerate(col):
            y = Inches(1.3 + ri * 1.35)
            add_rect(s, x, y, Inches(6.0), Inches(1.1), MID_BG)
            add_rect(s, x, y, Inches(0.55), Inches(1.1), ACCENT)
            add_textbox(s, num, x + Inches(0.05), y + Inches(0.2), Inches(0.5), Inches(0.7),
                        font_size=18, bold=True, color=HIGHLIGHT, align=PP_ALIGN.CENTER)
            add_textbox(s, title, x + Inches(0.65), y + Inches(0.12), Inches(5.0), Inches(0.55),
                        font_size=15, bold=True, color=WHITE)
            add_textbox(s, pg, x + Inches(0.65), y + Inches(0.65), Inches(5.0), Inches(0.3),
                        font_size=12, color=SUBTEXT)
    return s


def slide_dataset_overview(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    section_header_bar(s, "01  Dataset Overview — RACE")

    # Three stat boxes
    stats = [
        ("28,000+", "Reading Passages", HIGHLIGHT),
        ("100,000+", "MCQ Questions", GREEN),
        ("4 Options", "Per Question (A/B/C/D)", ORANGE),
    ]
    for i, (num, label, col) in enumerate(stats):
        x = Inches(0.4 + i * 4.25)
        add_rect(s, x, Inches(1.3), Inches(3.9), Inches(1.3), MID_BG)
        add_rect(s, x, Inches(1.3), Inches(3.9), Inches(0.08), col)
        add_textbox(s, num, x + Inches(0.15), Inches(1.42), Inches(3.6), Inches(0.65),
                    font_size=34, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_textbox(s, label, x + Inches(0.15), Inches(2.0), Inches(3.6), Inches(0.4),
                    font_size=14, color=SUBTEXT, align=PP_ALIGN.CENTER)

    # Left: dataset description
    txb = s.shapes.add_textbox(Inches(0.4), Inches(2.85), Inches(5.8), Inches(4.1))
    tf = txb.text_frame; tf.word_wrap = True
    add_para(tf, "What is RACE?", 18, bold=True, color=HIGHLIGHT)
    add_para(tf, "Reading Comprehension Assessment from Examinations — collected from Chinese middle and high school English exams.", 15, color=LIGHT_GREY, space_before=4)
    add_para(tf, "", 8)
    add_para(tf, "Raw CSV Columns", 16, bold=True, color=ORANGE)
    for col in ["article — full reading passage",
                "question — comprehension question",
                "A, B, C, D — four answer options",
                "answer — correct option label (A/B/C/D)"]:
        add_para(tf, f"  ▸  {col}", 14, color=LIGHT_GREY)

    # Right: splits table
    add_textbox(s, "Dataset Splits", Inches(6.6), Inches(2.85), Inches(6.3), Inches(0.45),
                font_size=16, bold=True, color=ORANGE)
    rows = [
        ("Split",      "Questions", "Expanded Rows\n(×4 options)", "Usage"),
        ("Train",      "70,280",    "281,120",                      "Model training"),
        ("Validation", "8,785",     "35,140",                       "Hyperparameter tuning"),
        ("Test",       "8,786",     "35,144",                       "Final evaluation"),
        ("TOTAL",      "87,851",    "351,404",                      "—"),
    ]
    col_w = [Inches(1.4), Inches(1.2), Inches(1.7), Inches(1.9)]
    col_x = [Inches(6.5), Inches(7.9), Inches(9.1), Inches(10.8)]
    for ri, row in enumerate(rows):
        y = Inches(3.4 + ri * 0.67)
        bg = ACCENT if ri == 0 else (MID_BG if ri % 2 == 1 else RGBColor(0x1E, 0x2D, 0x50))
        add_rect(s, Inches(6.5), y, Inches(6.4), Inches(0.62), bg)
        for ci, (cell, cw, cx) in enumerate(zip(row, col_w, col_x)):
            add_textbox(s, cell, cx + Inches(0.05), y + Inches(0.08), cw, Inches(0.5),
                        font_size=13 if ri > 0 else 13,
                        bold=(ri == 0 or ri == 5),
                        color=WHITE if ri == 0 else LIGHT_GREY)
    return s


def slide_dataset_expansion(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    section_header_bar(s, "01  Dataset — Per-Option Expansion & Two Processed CSVs")

    # Expansion diagram
    add_textbox(s, "Raw → Expanded Format", Inches(0.4), Inches(1.3), Inches(6), Inches(0.45),
                font_size=16, bold=True, color=ORANGE)

    boxes = [
        ("1 Question\n4 Options", Inches(0.4), Inches(1.9), Inches(2.4), Inches(1.3), ACCENT),
        ("×4\n(one row\nper option)", Inches(3.0), Inches(1.9), Inches(1.5), Inches(1.3), MID_BG),
        ("4 Rows\n1 positive\n3 negative", Inches(4.7), Inches(1.9), Inches(2.0), Inches(1.3), RGBColor(0x06, 0x4A, 0x2C)),
    ]
    for text, x, y, w, h, col in boxes:
        add_rect(s, x, y, w, h, col)
        add_textbox(s, text, x + Inches(0.1), y + Inches(0.1), w - Inches(0.2), h - Inches(0.2),
                    font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, "→", Inches(2.9), Inches(2.2), Inches(0.6), Inches(0.5),
                font_size=24, bold=True, color=HIGHLIGHT, align=PP_ALIGN.CENTER)
    add_textbox(s, "→", Inches(4.4), Inches(2.2), Inches(0.6), Inches(0.5),
                font_size=24, bold=True, color=HIGHLIGHT, align=PP_ALIGN.CENTER)

    add_textbox(s, "Class imbalance: 75% negative / 25% positive → requires class_weight='balanced'",
                Inches(0.4), Inches(3.35), Inches(6.5), Inches(0.45),
                font_size=13, italic=True, color=ORANGE)

    add_rect(s, 0, Inches(3.9), Inches(6.9), Inches(0.04), SUBTEXT)

    # Two processed CSVs
    add_textbox(s, "Two Processed CSVs Used in App", Inches(0.4), Inches(4.05), Inches(6.5), Inches(0.45),
                font_size=16, bold=True, color=ORANGE)

    csvs = [
        ("test.csv  (raw, per-question)",
         ["question_uid, id, article",
          "question, correct_answer, gold_answer_label",
          "distractor_1/2/3, distractor_label_1/2/3",
          "article_sentences"],
         "Used by lookup_sample_by_article() to fetch gold Q&A for RACE passages"),
        ("model_b_test_full.csv  (processed, per-option)",
         ["All raw columns +",
          "article_nopunct, question_nopunct",
          "correct_answer_nopunct",
          "TF-IDF feature arrays"],
         "Pre-computed features for faster Model B inference"),
    ]
    for i, (title, cols, note) in enumerate(csvs):
        x = Inches(0.4 + i * 3.4)
        add_rect(s, x, Inches(4.55), Inches(3.1), Inches(2.65), MID_BG)
        add_textbox(s, title, x + Inches(0.1), Inches(4.6), Inches(2.9), Inches(0.4),
                    font_size=13, bold=True, color=HIGHLIGHT)
        y_off = 5.1
        for col in cols:
            add_textbox(s, f"  • {col}", x + Inches(0.1), Inches(y_off), Inches(2.9), Inches(0.32),
                        font_size=12, color=LIGHT_GREY)
            y_off += 0.32
        add_textbox(s, note, x + Inches(0.1), Inches(y_off + 0.05), Inches(2.9), Inches(0.45),
                    font_size=11, italic=True, color=SUBTEXT)

    # Right: processed column listing
    txb = s.shapes.add_textbox(Inches(7.1), Inches(1.3), Inches(5.9), Inches(5.9))
    tf = txb.text_frame; tf.word_wrap = True
    add_para(tf, "test.csv — 16 Columns (8,786 rows)", 15, bold=True, color=ORANGE)
    cols16 = [
        "question_uid", "id", "article", "article_nopunct",
        "question", "question_nopunct", "correct_answer", "correct_answer_nopunct",
        "gold_answer_label", "distractor_1", "distractor_2", "distractor_3",
        "distractor_label_1", "distractor_label_2", "distractor_label_3",
        "article_sentences",
    ]
    for c in cols16:
        add_para(tf, f"  ▸  {c}", 13, color=LIGHT_GREY)
    return s


def slide_eda_distributions(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    section_header_bar(s, "02  EDA — Distributions & Class Balance")

    # Class balance
    add_textbox(s, "Label Distribution (per-option rows)", Inches(0.4), Inches(1.3), Inches(5.5), Inches(0.45),
                font_size=15, bold=True, color=ORANGE)

    bars = [("Negative (incorrect)", 75, RED), ("Positive (correct)", 25, GREEN)]
    for i, (label, pct, col) in enumerate(bars):
        y = Inches(1.85 + i * 0.85)
        add_rect(s, Inches(0.4), y, Inches(pct / 100 * 5.5), Inches(0.6), col)
        add_textbox(s, label, Inches(0.5), y + Inches(0.12), Inches(3.5), Inches(0.38),
                    font_size=14, bold=True, color=WHITE)
        add_textbox(s, f"{pct}%", Inches(5.0), y + Inches(0.12), Inches(0.8), Inches(0.38),
                    font_size=14, bold=True, color=col)

    add_textbox(s, "→  Severe imbalance (3:1) requires class_weight='balanced' in all classifiers",
                Inches(0.4), Inches(3.65), Inches(6.1), Inches(0.4),
                font_size=13, italic=True, color=ORANGE)

    add_rect(s, 0, Inches(4.15), Inches(6.9), Inches(0.04), SUBTEXT)

    # Answer label note — clean text, no chart
    add_textbox(s, "Answer Label Distribution", Inches(0.4), Inches(4.3), Inches(5.5), Inches(0.42),
                font_size=15, bold=True, color=ORANGE)
    txb2 = s.shapes.add_textbox(Inches(0.4), Inches(4.82), Inches(6.1), Inches(2.5))
    tf2 = txb2.text_frame; tf2.word_wrap = True
    add_para(tf2, "A: ~25.1%    B: ~24.9%    C: ~25.2%    D: ~24.8%", 15, bold=True, color=LIGHT_GREY)
    add_para(tf2, "", 6)
    add_para(tf2, "Labels are near-uniformly distributed across all four options.", 14, color=LIGHT_GREY, space_before=4)
    add_para(tf2, "This means the dataset has no answer-position bias — the model cannot", 14, color=LIGHT_GREY)
    add_para(tf2, "exploit a pattern like 'C is correct more often'. Any learned signal must", 14, color=LIGHT_GREY)
    add_para(tf2, "come from actual article/question/option content.", 14, color=LIGHT_GREY)
    add_para(tf2, "", 6)
    add_para(tf2, "No balancing was needed across labels A-D.", 13, italic=True, color=SUBTEXT)

    # Right: text length stats
    add_textbox(s, "Text Length Statistics", Inches(7.0), Inches(1.3), Inches(5.9), Inches(0.45),
                font_size=15, bold=True, color=ORANGE)
    rows = [
        ("Field", "Mean words", "Median", "Max"),
        ("Article", "~320", "~290", "~1,200"),
        ("Question", "~12", "~11", "~45"),
        ("Option (A/B/C/D)", "~8", "~7", "~50"),
        ("Correct Answer", "~8", "~7", "~50"),
    ]
    for ri, row in enumerate(rows):
        y = Inches(1.85 + ri * 0.72)
        bg = ACCENT if ri == 0 else (MID_BG if ri % 2 else RGBColor(0x1E, 0x2D, 0x50))
        add_rect(s, Inches(7.0), y, Inches(6.0), Inches(0.65), bg)
        xs = [Inches(7.05), Inches(9.0), Inches(10.3), Inches(11.4)]
        ws = [Inches(1.9), Inches(1.2), Inches(1.1), Inches(1.2)]
        for cell, x, w in zip(row, xs, ws):
            add_textbox(s, cell, x, y + Inches(0.1), w, Inches(0.45),
                        font_size=13, bold=(ri == 0), color=WHITE if ri == 0 else LIGHT_GREY)

    add_textbox(s, "Key insight: articles are ~40× longer than options → TF-IDF must handle long docs efficiently",
                Inches(7.0), Inches(5.7), Inches(6.0), Inches(0.5),
                font_size=13, italic=True, color=ORANGE)

    add_rect(s, Inches(6.9), Inches(1.15), Inches(0.04), Inches(6.2), SUBTEXT)
    return s


def slide_eda_template_wins(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    section_header_bar(s, "02  EDA — QGen Template Analysis & Vocabulary")

    # Template wins bar chart
    add_textbox(s, "Which template best matched the gold question? (train set wins)",
                Inches(0.4), Inches(1.3), Inches(7.5), Inches(0.45),
                font_size=15, bold=True, color=ORANGE)

    templates = [
        ("according_to_true",   14443, HIGHLIGHT),
        ("cloze_completion",    13663, GREEN),
        ("cloze_according_to",  10346, ORANGE),
        ("inference",            7445, RGBColor(0xAE, 0x81, 0xFF)),
        ("what_information",     6460, RED),
        ("main_idea_about",      5344, RGBColor(0xFF, 0xD1, 0x66)),
        ("main_idea_title",      4553, RGBColor(0x80, 0xED, 0x99)),
        ("not_true",             4315, RGBColor(0xFF, 0x87, 0xAB)),
        ("word_meaning",         2739, SUBTEXT),
    ]
    max_val = 14443
    for i, (name, count, col) in enumerate(templates):
        y = Inches(1.9 + i * 0.53)
        bar_w = count / max_val * Inches(5.5)
        add_rect(s, Inches(0.4), y, bar_w, Inches(0.42), col)
        add_textbox(s, name, Inches(0.5), y + Inches(0.05), Inches(3.8), Inches(0.35),
                    font_size=12, bold=True, color=DARK_BG if col != SUBTEXT else WHITE)
        add_textbox(s, f"{count:,}", Inches(0.4) + bar_w + Inches(0.1), y + Inches(0.05),
                    Inches(1.2), Inches(0.35), font_size=12, color=col)

    add_textbox(s, "→  'according_to_true' and 'cloze_completion' dominate — RACE exams heavily test\n    passage recall and sentence completion, not inference or title identification",
                Inches(0.4), Inches(6.9), Inches(6.7), Inches(0.5),
                font_size=12, italic=True, color=SUBTEXT)

    # Right: vocabulary — plain text, no boxes
    add_rect(s, Inches(6.9), Inches(1.15), Inches(0.04), Inches(6.2), SUBTEXT)

    txb_r = s.shapes.add_textbox(Inches(7.1), Inches(1.3), Inches(5.9), Inches(6.0))
    tf_r = txb_r.text_frame; tf_r.word_wrap = True
    add_para(tf_r, "Vocabulary Analysis", 16, bold=True, color=ORANGE)
    add_para(tf_r, "", 5)
    vocab_items = [
        ("Vocab size",         "50,000 tokens (max_features)"),
        ("n-gram range",       "(1, 2)  —  unigrams and bigrams"),
        ("min_df",             "5  —  ignores very rare terms"),
        ("max_df",             "0.95  —  ignores near-universal terms"),
        ("Stopwords removed",  "318 English stopwords (NLTK)"),
        ("Content-word rule",  "Token length > 2 AND not a stopword"),
    ]
    for key, val in vocab_items:
        add_para(tf_r, key, 14, bold=True, color=HIGHLIGHT, space_before=7)
        add_para(tf_r, f"  {val}", 13, color=LIGHT_GREY)
    return s


def slide_eda_cosine(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    section_header_bar(s, "02  EDA — Cosine Similarity Observations")

    add_textbox(s, "Three cosine similarity signals were computed for every (article, question, option) triple:",
                Inches(0.4), Inches(1.28), Inches(12.5), Inches(0.38), font_size=14, color=SUBTEXT)

    txb = s.shapes.add_textbox(Inches(0.4), Inches(1.8), Inches(12.5), Inches(5.5))
    tf = txb.text_frame; tf.word_wrap = True

    insights = [
        ("cos_article_option", HIGHLIGHT,
         "Correct options score significantly higher than incorrect ones (mean 0.35 vs 0.22).",
         "This is the single strongest individual predictor in the verifier feature set."),
        ("cos_question_option", GREEN,
         "Correct options align better with the question phrasing (mean 0.28 vs 0.18).",
         "Incorrect options often share vocabulary with the question but differ semantically."),
        ("cos_article_question", ORANGE,
         "Relatively uniform across correct and incorrect rows (~0.31 in both cases).",
         "The question is always drawn from the article, so this feature captures context, not correctness."),
    ]

    for feature, col, obs, interp in insights:
        add_para(tf, feature, 18, bold=True, color=col, space_before=14)
        add_para(tf, f"  Observation:   {obs}", 14, color=LIGHT_GREY)
        add_para(tf, f"  Interpretation: {interp}", 13, italic=True, color=SUBTEXT)

    add_para(tf, "", 8)
    add_para(tf, "All three signals were included in the 11-feature numeric block and scaled before concatenation with TF-IDF.", 13, italic=True, color=SUBTEXT)
    return s


def slide_eda_tfidf_vs_onehot(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    section_header_bar(s, "02  EDA — Why TF-IDF Instead of One-Hot Encoding?")

    add_textbox(s, "Five reasons TF-IDF was chosen as the primary text representation:",
                Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.38), font_size=14, color=SUBTEXT)

    reasons = [
        (HIGHLIGHT, "Main reason — Weighted word importance",
         "TF-IDF down-weights common words (appear in many articles) and up-weights rare but\n"
         "discriminative words. One-Hot treats every word as equally important — a critical flaw\n"
         "for reading comprehension where topic-specific terms matter most."),
        (ORANGE, "One-Hot would be enormous and useless",
         "Vocabulary across 28K articles is 80,000+ unique tokens. One-Hot produces a\n"
         "28K x 80K+ binary matrix — extremely sparse, no frequency or rarity information,\n"
         "and no benefit over TF-IDF which already handles sparsity better."),
        (GREEN, "Enabled cosine similarity features",
         "All 3 cosine similarity features (cos_article_option, cos_question_option,\n"
         "cos_article_question) are computed from TF-IDF vectors. One-Hot cosines are\n"
         "dominated by stopword matches and give near-meaningless similarity scores."),
        (RGBColor(0xAE, 0x81, 0xFF), "Better fit for LinearSVC and Logistic Regression",
         "LinearSVC and LR work best with real-valued, normalized features. TF-IDF produces\n"
         "L2-normalized continuous values per token — ideal for these linear classifiers.\n"
         "One-Hot binary inputs remove the gradient signal that TF-IDF preserves."),
        (RED, "Also used for Question Generation (QGen)",
         "The assignment notes that candidate sentences for QGen can be extracted using\n"
         "One-Hot keyword overlap, but TF-IDF keyword overlap was used as the alternative.\n"
         "TF-IDF sentence scoring (cosine to correct answer) picks more relevant source sentences."),
    ]

    for i, (col, title, body) in enumerate(reasons):
        y = Inches(1.75 + i * 1.12)
        add_rect(s, Inches(0.3), y, Inches(12.7), Inches(1.02), MID_BG)
        add_rect(s, Inches(0.3), y, Inches(0.14), Inches(1.02), col)
        add_textbox(s, title, Inches(0.58), y + Inches(0.06), Inches(12.0), Inches(0.36),
                    font_size=15, bold=True, color=col)
        add_textbox(s, body, Inches(0.58), y + Inches(0.44), Inches(12.0), Inches(0.54),
                    font_size=13, color=LIGHT_GREY)
    return s


def slide_statistical_analysis(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    section_header_bar(s, "03  Statistical Analysis")

    # Left: unsupervised experiments
    add_textbox(s, "Unsupervised Clustering — K-Means vs GMM",
                Inches(0.4), Inches(1.3), Inches(6.2), Inches(0.45),
                font_size=15, bold=True, color=ORANGE)

    rows = [
        ("Method",       "Features",     "Silhouette", "Purity"),
        ("K-Means k=2",  "TF-IDF",       "0.0062",     "0.75"),
        ("K-Means k=4",  "TF-IDF",       "0.0051",     "0.75"),
        ("K-Means k=6",  "TF-IDF",       "0.0051",     "0.75"),
        ("K-Means k=8",  "TF-IDF",       "0.0051",     "0.75"),
        ("GMM (4 comp)", "11 Numeric",   "0.0735 ✓",   "0.75"),
    ]
    for ri, row in enumerate(rows):
        y = Inches(1.9 + ri * 0.65)
        bg = ACCENT if ri == 0 else (RGBColor(0x1E, 0x2D, 0x50) if ri % 2 == 0 else MID_BG)
        if ri == 5: bg = RGBColor(0x06, 0x3A, 0x2C)
        add_rect(s, Inches(0.4), y, Inches(6.2), Inches(0.6), bg)
        xs = [Inches(0.5), Inches(2.3), Inches(3.9), Inches(5.4)]
        ws = [Inches(1.8), Inches(1.6), Inches(1.5), Inches(1.1)]
        for cell, x, w in zip(row, xs, ws):
            add_textbox(s, cell, x, y + Inches(0.1), w, Inches(0.42),
                        font_size=13, bold=(ri == 0 or ri == 5),
                        color=WHITE if ri == 0 else (GREEN if ri == 5 else LIGHT_GREY))

    add_textbox(s,
        "Purity is constant at 0.75 = majority class proportion.\n"
        "Silhouette tells the true story: GMM on numeric features (0.0735)\n"
        "forms more coherent clusters than K-Means on sparse TF-IDF (0.006).",
        Inches(0.4), Inches(6.2), Inches(6.2), Inches(1.1),
        font_size=13, italic=True, color=SUBTEXT)

    add_rect(s, Inches(6.9), Inches(1.15), Inches(0.04), Inches(6.2), SUBTEXT)

    # Right: Semi-supervised + BIC/AIC
    add_textbox(s, "Semi-supervised: Label Propagation (5% labels)",
                Inches(7.1), Inches(1.3), Inches(6.0), Inches(0.45),
                font_size=15, bold=True, color=ORANGE)

    rows2 = [
        ("Model",              "Val Acc",  "F1-macro", "F1-correct"),
        ("Supervised LR\n(100% labels)", "0.5312", "—",     "0.3746"),
        ("Label Prop\n(5% labels)",      "0.7357", "0.4465","0.0464 ⚠"),
    ]
    for ri, row in enumerate(rows2):
        y = Inches(1.9 + ri * 0.82)
        bg = ACCENT if ri == 0 else (MID_BG if ri == 1 else RGBColor(0x3A, 0x10, 0x10))
        add_rect(s, Inches(7.1), y, Inches(6.0), Inches(0.72), bg)
        xs = [Inches(7.2), Inches(9.2), Inches(10.3), Inches(11.3)]
        ws = [Inches(2.0), Inches(1.1), Inches(1.0), Inches(1.1)]
        for cell, x, w in zip(row, xs, ws):
            add_textbox(s, cell, x, y + Inches(0.08), w, Inches(0.58),
                        font_size=13, bold=(ri == 0), color=WHITE if ri == 0 else LIGHT_GREY)

    add_textbox(s,
        "⚠  LP achieves higher overall accuracy by collapsing to the\n"
        "majority negative class. F1 for the 'correct' class drops from\n"
        "0.375 → 0.046. This is a known failure mode on imbalanced\n"
        "data with graph-based label propagation.",
        Inches(7.1), Inches(3.62), Inches(6.0), Inches(1.3),
        font_size=13, italic=True, color=ORANGE)

    add_rect(s, Inches(7.1), Inches(5.05), Inches(6.0), Inches(0.04), SUBTEXT)

    txb_bic = s.shapes.add_textbox(Inches(7.1), Inches(5.15), Inches(6.0), Inches(2.22))
    tf_bic = txb_bic.text_frame; tf_bic.word_wrap = True
    add_para(tf_bic, "GMM Model Selection — BIC & AIC", 15, bold=True, color=ORANGE)
    add_para(tf_bic, "", 4)
    add_para(tf_bic,
        "BIC (Bayesian Information Criterion) and AIC (Akaike Information Criterion) measure how "
        "well a model fits the data while penalising model complexity. Lower values are better.",
        13, color=LIGHT_GREY)
    add_para(tf_bic, "", 5)
    add_para(tf_bic, "  BIC:                    -641,640.7", 13, bold=True, color=HIGHLIGHT)
    add_para(tf_bic, "  AIC:                    -644,920.6", 13, bold=True, color=HIGHLIGHT)
    add_para(tf_bic, "  Log-likelihood/sample:       1.148", 13, color=LIGHT_GREY)
    add_para(tf_bic, "", 5)
    add_para(tf_bic,
        "Strongly negative BIC/AIC confirms the 4-component GMM fits the 11 numeric "
        "features well without overfitting.",
        12, italic=True, color=SUBTEXT)
    return s


def slide_stat_pseudo_labels(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    section_header_bar(s, "03  Statistical Analysis — Pseudo-Label Quality Scoring")

    add_textbox(s,
        "QGen candidates are pseudo-labeled using a weighted composite score against the gold question:",
        Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.4), font_size=15, color=LIGHT_GREY)

    # Formula box
    add_rect(s, Inches(1.0), Inches(1.85), Inches(11.2), Inches(1.1), MID_BG)
    add_rect(s, Inches(1.0), Inches(1.85), Inches(0.12), Inches(1.1), HIGHLIGHT)
    add_textbox(s,
        "quality_score  =  0.30 × BLEU-1  +  0.20 × BLEU-2  +  0.25 × ROUGE-L  +  0.25 × METEOR",
        Inches(1.2), Inches(2.0), Inches(10.8), Inches(0.65),
        font_size=20, bold=True, color=HIGHLIGHT, align=PP_ALIGN.CENTER)

    txb = s.shapes.add_textbox(Inches(0.4), Inches(3.15), Inches(12.5), Inches(4.2))
    tf = txb.text_frame; tf.word_wrap = True
    add_para(tf, "Why this weighting?", 16, bold=True, color=ORANGE)
    reasons = [
        "BLEU-1 (30%): Unigram precision — most important for capturing key content words",
        "BLEU-2 (20%): Bigram precision — rewards preserving phrase structure",
        "ROUGE-L (25%): Longest common subsequence — rewards word order similarity",
        "METEOR (25%): Synonym + stem aware — penalizes exact-match bias in BLEU",
    ]
    for r in reasons:
        add_para(tf, f"  ▸  {r}", 14, color=LIGHT_GREY, space_before=5)

    add_para(tf, "", 8)
    add_para(tf, "Statistical outcome:", 16, bold=True, color=ORANGE)
    add_para(tf, "  ▸  Train candidates: 1,132,145 rows → positive rate: 6.21% (best candidate per question)", 14, color=LIGHT_GREY)
    add_para(tf, "  ▸  The ranker learns to distinguish high-quality generated questions from low-quality ones", 14, color=LIGHT_GREY)
    add_para(tf, "  ▸  This is the bridge between unsupervised template generation and supervised ranking", 14, color=LIGHT_GREY)
    return s


def slide_preprocessing(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    section_header_bar(s, "04  Data Preprocessing Pipeline")

    steps = [
        ("1", "Text Normalisation",
         "lowercase → collapse whitespace → strip leading/trailing spaces\n"
         "Applied to: article, question, all options",
         HIGHLIGHT),
        ("2", "Punctuation Removal",
         "str.translate(str.maketrans('', '', string.punctuation))\n"
         "Used for overlap feature computation; TF-IDF sees original text",
         GREEN),
        ("3", "Tokenisation",
         "regex r'\\b[a-zA-Z][a-zA-Z0-9]*\\b' → word token list\n"
         "Removes numbers-only tokens; preserves alphanumeric",
         ORANGE),
        ("4", "Stopword Removal",
         "NLTK English stopwords (318 words) + token length > 2\n"
         "Applied only for content-word overlap features, not TF-IDF",
         RGBColor(0xAE, 0x81, 0xFF)),
        ("5", "TF-IDF Vectorisation",
         "TfidfVectorizer(max_features=50000, min_df=5, max_df=0.95, ngram_range=(1,2))\n"
         "Fit on train set only → transform val/test to prevent data leakage",
         RED),
        ("6", "Numeric Feature Scaling",
         "StandardScaler fit on train numeric features → transform val/test\n"
         "Ensures cosine/overlap features are on the same scale as each other",
         ORANGE),
        ("7", "Feature Concatenation",
         "scipy.sparse.hstack([tfidf_matrix, scaled_numeric_matrix])\n"
         "Produces 50,011-dim combined feature matrix per row",
         HIGHLIGHT),
    ]

    for i, (num, title, detail, col) in enumerate(steps):
        col_i = i % 2
        row_i = i // 2
        x = Inches(0.35 + col_i * 6.55)
        y = Inches(1.3 + row_i * 1.5)
        if i == 6:
            x = Inches(3.45); y = Inches(7.3 - 1.5)  # centred last step

        add_rect(s, x, y, Inches(6.3), Inches(1.35), MID_BG)
        add_rect(s, x, y, Inches(0.45), Inches(1.35), col)
        add_textbox(s, num, x + Inches(0.0), y + Inches(0.35), Inches(0.45), Inches(0.55),
                    font_size=18, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        add_textbox(s, title, x + Inches(0.55), y + Inches(0.08), Inches(5.6), Inches(0.38),
                    font_size=14, bold=True, color=col)
        add_textbox(s, detail, x + Inches(0.55), y + Inches(0.5), Inches(5.6), Inches(0.78),
                    font_size=12, color=LIGHT_GREY)
    return s


def slide_preprocessing_features(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    section_header_bar(s, "04  Feature Engineering — 11 Numeric Features (Model A Verifier)")

    add_textbox(s, "These 11 features are computed for every (article, question, option) triple and concatenated with TF-IDF:",
                Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.4), font_size=14, color=LIGHT_GREY)

    features = [
        ("cos_article_option",          "Cosine similarity (TF-IDF vectors) between article and option",          "Strongest single predictor", HIGHLIGHT),
        ("cos_question_option",         "Cosine similarity between question and option",                           "Second strongest predictor", HIGHLIGHT),
        ("cos_article_question",        "Cosine similarity between article and question",                          "Context relevance", HIGHLIGHT),
        ("article_len",                 "Token count of the article",                                              "Longer articles → harder", GREEN),
        ("question_len",                "Token count of the question",                                             "Question complexity", GREEN),
        ("option_len",                  "Token count of the option",                                               "Answer verbosity", GREEN),
        ("question_option_overlap",     "Count of shared content words between question and option",               "Direct matching signal", ORANGE),
        ("article_option_overlap",      "Count of shared content words between article and option",                "Passage evidence", ORANGE),
        ("article_question_overlap",    "Count of shared content words between article and question",              "Question grounding", ORANGE),
        ("question_option_overlap_ratio","Jaccard ratio of question ∩ option / question ∪ option",                "Normalised for length", RED),
        ("article_option_overlap_ratio", "Jaccard ratio of article ∩ option / article ∪ option",                  "Normalised for length", RED),
    ]

    col_groups = [features[:6], features[6:]]
    for ci, group in enumerate(col_groups):
        x = Inches(0.3 + ci * 6.55)
        for ri, (name, desc, note, col) in enumerate(group):
            y = Inches(1.85 + ri * 0.86)
            add_rect(s, x, y, Inches(6.3), Inches(0.8), MID_BG)
            add_rect(s, x, y, Inches(0.12), Inches(0.8), col)
            add_textbox(s, name, x + Inches(0.2), y + Inches(0.04), Inches(5.9), Inches(0.34),
                        font_size=13, bold=True, color=col)
            add_textbox(s, desc, x + Inches(0.2), y + Inches(0.4), Inches(4.5), Inches(0.34),
                        font_size=12, color=LIGHT_GREY)
            add_textbox(s, note, x + Inches(4.75), y + Inches(0.4), Inches(1.45), Inches(0.34),
                        font_size=11, italic=True, color=SUBTEXT)
    return s


def slide_model_a_selection(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    section_header_bar(s, "05  Model A — Selection Process (Answer Verifier)")

    add_textbox(s, "Seven configurations evaluated on the validation set (option-level F1):",
                Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.38), font_size=14, color=LIGHT_GREY)

    rows = [
        ("Model Configuration",                "Features",      "Val F1",  "Selected?"),
        ("Naive Bayes",                         "TF-IDF only",   "0.4286",  ""),
        ("Random Forest",                       "Numeric only",  "0.4632",  ""),
        ("Logistic Regression",                 "TF-IDF only",   "0.4945",  ""),
        ("Hard Voting Ensemble (LR+SVM+NB)",    "TF-IDF only",   "0.5151",  ""),
        ("LinearSVC",                           "TF-IDF only",   "0.5031",  ""),
        ("Logistic Regression",                 "Combined",      "0.5081",  ""),
        ("LinearSVC",                           "Combined",      "0.5176",  "✅ BEST"),
    ]
    col_xs = [Inches(0.35), Inches(5.7), Inches(9.2), Inches(11.0)]
    col_ws = [Inches(5.3), Inches(3.4), Inches(1.7), Inches(2.0)]
    for ri, row in enumerate(rows):
        y = Inches(1.75 + ri * 0.62)
        is_best = ri == 7
        bg = ACCENT if ri == 0 else (RGBColor(0x06, 0x4A, 0x2C) if is_best else (MID_BG if ri % 2 else RGBColor(0x1E, 0x2D, 0x50)))
        add_rect(s, Inches(0.35), y, Inches(12.6), Inches(0.57), bg)
        for cell, x, w in zip(row, col_xs, col_ws):
            add_textbox(s, cell, x + Inches(0.08), y + Inches(0.08), w, Inches(0.42),
                        font_size=13 if not is_best else 14,
                        bold=(ri == 0 or is_best),
                        color=GREEN if is_best and cell == "✅ BEST" else (WHITE if ri == 0 else LIGHT_GREY))

    txb = s.shapes.add_textbox(Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.72))
    tf = txb.text_frame; tf.word_wrap = True
    add_para(tf, "Why LinearSVC + Combined wins:", 15, bold=True, color=ORANGE)
    add_para(tf, "  ▸  LinearSVC handles high-dimensional sparse TF-IDF better than tree models (no curse of dimensionality for linear boundaries)", 13, color=LIGHT_GREY)
    add_para(tf, "  ▸  Combined features add cosine similarity context that TF-IDF alone cannot capture for short option texts", 13, color=LIGHT_GREY)
    return s


def slide_model_a_training(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    section_header_bar(s, "05  Model A — Training Process")

    # Left: verifier training
    add_textbox(s, "Answer Verifier Training", Inches(0.4), Inches(1.3), Inches(6.0), Inches(0.42),
                font_size=16, bold=True, color=ORANGE)

    steps_v = [
        ("Input", "281,120 rows (train) — article + question + option"),
        ("TF-IDF", "Fit TfidfVectorizer on train set only → 50K-dim sparse matrix"),
        ("Numeric", "Compute 11 features per row → fit StandardScaler on train"),
        ("Combine", "scipy.sparse.hstack([tfidf_train, scaled_numeric_train])"),
        ("Train", "LinearSVC(class_weight='balanced') — handles 3:1 imbalance"),
        ("Validate", "Evaluate on 35,140 val rows → F1 = 0.5176"),
        ("Test", "Final evaluation on 35,144 test rows → Question acc = 37.3%"),
    ]
    for i, (label, detail) in enumerate(steps_v):
        y = Inches(1.85 + i * 0.7)
        add_rect(s, Inches(0.35), y, Inches(6.2), Inches(0.62), MID_BG)
        add_rect(s, Inches(0.35), y, Inches(0.85), Inches(0.62), ACCENT)
        add_textbox(s, label, Inches(0.38), y + Inches(0.1), Inches(0.8), Inches(0.42),
                    font_size=13, bold=True, color=HIGHLIGHT, align=PP_ALIGN.CENTER)
        add_textbox(s, detail, Inches(1.28), y + Inches(0.1), Inches(5.1), Inches(0.42),
                    font_size=13, color=LIGHT_GREY)

    add_rect(s, Inches(6.75), Inches(1.15), Inches(0.04), Inches(6.2), SUBTEXT)

    # Right: QGen training
    add_textbox(s, "QGen Ranker Training", Inches(7.0), Inches(1.3), Inches(6.0), Inches(0.42),
                font_size=16, bold=True, color=ORANGE)

    steps_q = [
        ("Sentences", "Split each article into sentences (SpaCy / regex)"),
        ("Top-2", "Select top-2 sentences by TF-IDF cosine to correct answer"),
        ("Templates", "Generate 4–9 candidate questions per sentence via templates"),
        ("Score", "Compute quality_score vs gold question (BLEU+ROUGE+METEOR)"),
        ("Label", "Mark best candidate per question as positive (rate=6.21%)"),
        ("Features", "Build 27-dim feature vector per candidate"),
        ("Train", "LinearSVC ranker on 1,132,145 candidates (train set)"),
        ("Eval", "Pick top-1 ranked candidate → compare to gold question"),
    ]
    for i, (label, detail) in enumerate(steps_q):
        y = Inches(1.85 + i * 0.67)
        add_rect(s, Inches(7.0), y, Inches(6.0), Inches(0.6), MID_BG)
        add_rect(s, Inches(7.0), y, Inches(0.85), Inches(0.6), ACCENT)
        add_textbox(s, label, Inches(7.03), y + Inches(0.1), Inches(0.8), Inches(0.4),
                    font_size=12, bold=True, color=HIGHLIGHT, align=PP_ALIGN.CENTER)
        add_textbox(s, detail, Inches(7.93), y + Inches(0.1), Inches(5.0), Inches(0.4),
                    font_size=12, color=LIGHT_GREY)
    return s


def slide_model_a_qgen_features(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    section_header_bar(s, "05  Model A — 27 QGen Features")

    groups = [
        ("Cosine Similarities (TF-IDF)", HIGHLIGHT, [
            "cos_q_sent — question vs source sentence",
            "cos_q_ans  — question vs correct answer",
            "cos_q_ctx  — question vs full article context",
        ]),
        ("Lexical Scores", GREEN, [
            "sentence_score      — TF-IDF relevance of sentence to article",
            "contains_answer     — 1 if source sentence contains answer text",
            "sentence_position   — position of sentence in article (0=first)",
        ]),
        ("Length Features", ORANGE, [
            "generated_q_len     — token count of generated question",
            "source_sentence_len — token count of source sentence",
            "correct_answer_len  — token count of correct answer",
        ]),
        ("Overlap Features", RGBColor(0xAE, 0x81, 0xFF), [
            "q_sent_overlap / q_ans_overlap        — content word counts",
            "q_sent_overlap_ratio / q_ans_overlap_ratio  — Jaccard ratios",
        ]),
        ("Wh-type One-hot (7)", RED, [
            "wh_who, wh_what, wh_where, wh_when, wh_why, wh_how, wh_which",
        ]),
        ("Answer Type One-hot (4)", RGBColor(0xFF, 0xD1, 0x66), [
            "answer_type_who, answer_type_what, answer_type_where, answer_type_when",
        ]),
        ("Template Group One-hot (3)", SUBTEXT, [
            "template_is_cloze, template_is_according, template_is_generic",
        ]),
    ]

    col_data = [groups[:4], groups[4:]]
    for ci, col_groups in enumerate(col_data):
        x = Inches(0.3 + ci * 6.5)
        y_off = 1.3
        for group_name, col, items in col_groups:
            add_rect(s, x, Inches(y_off), Inches(6.25), Inches(0.36), col)
            add_textbox(s, group_name, x + Inches(0.1), Inches(y_off + 0.04),
                        Inches(6.0), Inches(0.3), font_size=13, bold=True, color=DARK_BG)
            y_off += 0.38
            for item in items:
                add_rect(s, x, Inches(y_off), Inches(6.25), Inches(0.4), MID_BG)
                add_textbox(s, f"  {item}", x + Inches(0.05), Inches(y_off + 0.05),
                            Inches(6.1), Inches(0.32), font_size=12, color=LIGHT_GREY)
                y_off += 0.42
            y_off += 0.12
    return s


def slide_model_a_results(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    section_header_bar(s, "05  Model A — Evaluation Results (Test Set)")

    # Verifier results
    add_textbox(s, "Answer Verifier — LinearSVC + Combined", Inches(0.4), Inches(1.3), Inches(6.0), Inches(0.42),
                font_size=16, bold=True, color=ORANGE)

    metrics_v = [
        ("Question-level Accuracy", "0.3732", "37.3% vs 25% random chance", HIGHLIGHT),
        ("Option-level F1 (macro)", "0.5176", "Balanced class performance", GREEN),
        ("BLEU-1",  "0.4806", "", LIGHT_GREY),
        ("BLEU-4",  "0.3325", "", LIGHT_GREY),
        ("ROUGE-1", "0.5043", "", LIGHT_GREY),
        ("ROUGE-L", "0.4971", "", LIGHT_GREY),
        ("METEOR",  "0.4400", "", LIGHT_GREY),
    ]
    for i, (name, val, note, col) in enumerate(metrics_v):
        y = Inches(1.85 + i * 0.67)
        add_rect(s, Inches(0.35), y, Inches(6.2), Inches(0.6), MID_BG)
        add_textbox(s, name, Inches(0.45), y + Inches(0.1), Inches(3.5), Inches(0.42),
                    font_size=13, color=LIGHT_GREY)
        add_textbox(s, val, Inches(4.0), y + Inches(0.1), Inches(1.2), Inches(0.42),
                    font_size=16, bold=True, color=col, align=PP_ALIGN.CENTER)
        if note:
            add_textbox(s, note, Inches(5.3), y + Inches(0.1), Inches(1.1), Inches(0.42),
                        font_size=11, italic=True, color=SUBTEXT)

    add_rect(s, Inches(6.75), Inches(1.15), Inches(0.04), Inches(6.2), SUBTEXT)

    # QGen results
    add_textbox(s, "Question Generator — LinearSVC Ranker", Inches(7.0), Inches(1.3), Inches(6.0), Inches(0.42),
                font_size=16, bold=True, color=ORANGE)

    metrics_rows = [
        ("Metric",   "Baseline\n(no ranker)", "LinearSVC\nRanker", "RF\nRanker"),
        ("BLEU-1",   "0.1185",               "0.1811",             "—"),
        ("BLEU-2",   "0.0472",               "0.1018",             "—"),
        ("BLEU-4",   "0.0238",               "0.0607",             "—"),
        ("ROUGE-1",  "0.1906",               "0.2362",             "—"),
        ("ROUGE-2",  "0.0405",               "0.0927",             "—"),
        ("ROUGE-L",  "0.1700",               "0.2091",             "—"),
        ("METEOR",   "0.0914",               "0.1598",             "—"),
    ]
    xs = [Inches(7.05), Inches(9.3), Inches(10.8), Inches(12.1)]
    ws = [Inches(2.2), Inches(1.5), Inches(1.3), Inches(1.0)]
    for ri, row in enumerate(metrics_rows):
        y = Inches(1.85 + ri * 0.63)
        bg = ACCENT if ri == 0 else (MID_BG if ri % 2 else RGBColor(0x1E, 0x2D, 0x50))
        if ri > 0 and float(row[2]) == max(float(r[2]) for r in metrics_rows[1:] if r[2] != "—"):
            bg = RGBColor(0x06, 0x4A, 0x2C)
        add_rect(s, Inches(7.0), y, Inches(6.1), Inches(0.58), bg)
        for cell, x, w in zip(row, xs, ws):
            add_textbox(s, cell, x, y + Inches(0.06), w, Inches(0.48),
                        font_size=13, bold=(ri == 0), color=WHITE if ri == 0 else LIGHT_GREY)

    add_textbox(s, "+52.8% METEOR vs baseline — ranker meaningfully improves template selection",
                Inches(7.0), Inches(6.95), Inches(6.1), Inches(0.4),
                font_size=13, italic=True, color=GREEN)
    return s


def slide_model_b_overview(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    section_header_bar(s, "06  Model B — Distractor & Hint Generator Overview")

    txb = s.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.8))
    tf = txb.text_frame; tf.word_wrap = True
    add_para(tf, "Model B operates after Model A — given (article, question, correct_answer), it generates:", 15, color=LIGHT_GREY)

    components = [
        ("Distractor Generator",
         "3 plausible but incorrect options (A/B/C/D minus correct). Candidate pool: top-2 article\n"
         "sentences × ~9 templates ≈ 18 candidates per question. Random Forest ranker selects top-3.",
         ORANGE),
        ("Hint Scorer",
         "3 graduated hints from article sentences, ordered from least to most revealing.\n"
         "Logistic Regression ranks all article sentences; top-3 by score become the hints.",
         HIGHLIGHT),
    ]

    for i, (title, body, col) in enumerate(components):
        y = Inches(2.0 + i * 2.2)
        add_rect(s, Inches(0.3), y, Inches(12.7), Inches(1.9), MID_BG)
        add_rect(s, Inches(0.3), y, Inches(0.14), Inches(1.9), col)
        add_textbox(s, title, Inches(0.6), y + Inches(0.1), Inches(12.0), Inches(0.5),
                    font_size=18, bold=True, color=col)
        add_textbox(s, body, Inches(0.6), y + Inches(0.6), Inches(12.0), Inches(1.15),
                    font_size=15, color=LIGHT_GREY)

    add_textbox(s, "Shared TF-IDF vectorizer (50K vocab) — same parameters as Model A, fit on per-question corpus",
                Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.38),
                font_size=14, italic=True, color=SUBTEXT)

    add_rect(s, Inches(0.3), Inches(6.48), Inches(12.7), Inches(0.04), SUBTEXT)

    # Candidate scale
    cols_info = [
        ("1,288,390", "Train candidates\n(distractor)", ORANGE),
        ("70,280", "Train questions", GREEN),
        ("~18", "Candidates per\nquestion", HIGHLIGHT),
        ("6.21%", "Positive rate\n(1 of 18 is gold)", RED),
    ]
    for i, (num, label, col) in enumerate(cols_info):
        x = Inches(0.5 + i * 3.2)
        add_rect(s, x, Inches(6.6), Inches(2.9), Inches(0.82), RGBColor(0x0A, 0x10, 0x20))
        add_textbox(s, num, x + Inches(0.1), Inches(6.62), Inches(2.7), Inches(0.45),
                    font_size=22, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_textbox(s, label, x + Inches(0.1), Inches(7.04), Inches(2.7), Inches(0.35),
                    font_size=11, color=SUBTEXT, align=PP_ALIGN.CENTER)
    return s


def slide_model_b_features(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    section_header_bar(s, "06  Model B — Feature Sets")

    # 14 Distractor features
    add_textbox(s, "14 Distractor Features", Inches(0.4), Inches(1.3), Inches(6.0), Inches(0.42),
                font_size=16, bold=True, color=ORANGE)

    d_features = [
        ("cos_candidate_answer",    "Cosine sim: candidate ↔ correct answer",  HIGHLIGHT),
        ("cos_candidate_question",  "Cosine sim: candidate ↔ question",         HIGHLIGHT),
        ("cos_candidate_article",   "Cosine sim: candidate ↔ article",          HIGHLIGHT),
        ("candidate_word_len",      "Token count of candidate phrase",           GREEN),
        ("answer_word_len",         "Token count of correct answer",             GREEN),
        ("length_diff",             "Absolute token count difference",           GREEN),
        ("char_overlap_answer",     "Character-level Jaccard with answer",       ORANGE),
        ("word_overlap_answer",     "Content-word Jaccard with answer",          ORANGE),
        ("word_overlap_question",   "Content-word Jaccard with question",        ORANGE),
        ("candidate_frequency_article","How often candidate appears in article", RED),
        ("candidate_in_question",   "1 if candidate appears in question text",  RED),
        ("candidate_same_as_answer","1 if candidate == correct answer (filter)", RED),
        ("candidate_contains_answer","1 if candidate contains answer as substring", RED),
        ("answer_contains_candidate","1 if answer contains candidate as substring", RED),
    ]
    for i, (name, desc, col) in enumerate(d_features):
        y = Inches(1.85 + i * 0.35)
        add_rect(s, Inches(0.35), y, Inches(6.2), Inches(0.31), MID_BG)
        add_rect(s, Inches(0.35), y, Inches(0.08), Inches(0.31), col)
        add_textbox(s, name, Inches(0.5), y + Inches(0.03), Inches(2.6), Inches(0.26),
                    font_size=11, bold=True, color=col)
        add_textbox(s, desc, Inches(3.15), y + Inches(0.03), Inches(3.35), Inches(0.26),
                    font_size=11, color=LIGHT_GREY)

    add_rect(s, Inches(6.75), Inches(1.15), Inches(0.04), Inches(6.2), SUBTEXT)

    # 8 Hint features
    add_textbox(s, "8 Hint Features", Inches(7.0), Inches(1.3), Inches(6.0), Inches(0.42),
                font_size=16, bold=True, color=ORANGE)

    h_features = [
        ("cos_question_sentence",         "Cosine sim: question ↔ sentence",         HIGHLIGHT),
        ("cos_answer_sentence",           "Cosine sim: answer ↔ sentence",            HIGHLIGHT),
        ("cos_question_answer_sentence",  "Cosine sim: (q+a) combined ↔ sentence",   HIGHLIGHT),
        ("question_sentence_overlap",     "Content-word count: question ∩ sentence",  ORANGE),
        ("answer_sentence_overlap",       "Content-word count: answer ∩ sentence",    ORANGE),
        ("sentence_position_norm",        "Normalised position (0=first, 1=last)",    GREEN),
        ("sentence_word_len",             "Token count of the sentence",              GREEN),
        ("contains_correct_answer",       "1 if sentence contains answer text",       RED),
    ]
    for i, (name, desc, col) in enumerate(h_features):
        y = Inches(1.85 + i * 0.6)
        add_rect(s, Inches(7.0), y, Inches(6.0), Inches(0.54), MID_BG)
        add_rect(s, Inches(7.0), y, Inches(0.08), Inches(0.54), col)
        add_textbox(s, name, Inches(7.15), y + Inches(0.04), Inches(2.8), Inches(0.28),
                    font_size=12, bold=True, color=col)
        add_textbox(s, desc, Inches(7.15), y + Inches(0.28), Inches(5.75), Inches(0.24),
                    font_size=12, color=LIGHT_GREY)
    return s


def slide_model_b_training(prs):
    """Slide 18 — distractor model selection (clean, one focus)."""
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    section_header_bar(s, "06  Model B — Distractor Ranker: Model Selection")

    add_textbox(s, "Two classifiers compared on the validation set (5,000 questions, top-3 distractors each):",
                Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.38), font_size=14, color=SUBTEXT)

    # Model comparison table — wide & readable
    rows = [
        ("Model",                "BLEU-1", "BLEU-2", "BLEU-4", "ROUGE-1", "ROUGE-L", "METEOR"),
        ("Logistic Regression",  "0.7582", "0.7208", "0.6476", "0.7631",  "0.7631",  "0.7276"),
        ("Random Forest (n=200)","0.7993", "0.7499", "0.6652", "0.8035",  "0.8035",  "0.7588"),
    ]
    col_xs = [Inches(0.35), Inches(4.2), Inches(5.75), Inches(7.3), Inches(8.85), Inches(10.4), Inches(11.85)]
    col_ws = [Inches(3.8), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.45), Inches(1.45)]
    for ri, row in enumerate(rows):
        y = Inches(1.75 + ri * 0.88)
        is_rf = ri == 2
        bg = ACCENT if ri == 0 else (RGBColor(0x06, 0x4A, 0x2C) if is_rf else MID_BG)
        add_rect(s, Inches(0.35), y, Inches(12.9), Inches(0.78), bg)
        for cell, x, w in zip(row, col_xs, col_ws):
            add_textbox(s, cell, x + Inches(0.05), y + Inches(0.14), w, Inches(0.5),
                        font_size=14, bold=(ri == 0 or is_rf),
                        color=GREEN if is_rf and cell not in [row[0]] else (WHITE if ri == 0 else LIGHT_GREY))

    # Winner callout
    add_rect(s, Inches(0.35), Inches(4.4), Inches(12.9), Inches(1.05), RGBColor(0x06, 0x4A, 0x2C))
    add_rect(s, Inches(0.35), Inches(4.4), Inches(0.14), Inches(1.05), GREEN)
    add_textbox(s, "Random Forest selected  (every metric higher)",
                Inches(0.65), Inches(4.48), Inches(12.3), Inches(0.42),
                font_size=18, bold=True, color=GREEN)
    add_textbox(s,
        "RF captures non-linear interactions between cosine similarity, word length, and candidate frequency "
        "that a linear boundary (LR) cannot separate. n_estimators=200, class_weight='balanced'.",
        Inches(0.65), Inches(4.9), Inches(12.3), Inches(0.48),
        font_size=14, color=LIGHT_GREY)

    add_rect(s, 0, Inches(5.65), SLIDE_W, Inches(0.05), SUBTEXT)

    # Training setup summary — 3 clean boxes
    add_textbox(s, "Training Setup", Inches(0.35), Inches(5.78), Inches(4.0), Inches(0.38),
                font_size=15, bold=True, color=ORANGE)

    boxes = [
        ("1,288,390", "candidate rows\n(train set)", ORANGE),
        ("6.21%", "positive rate\n(1 gold per ~18 candidates)", GREEN),
        ("14", "features per\ncandidate row", HIGHLIGHT),
    ]
    for i, (num, label, col) in enumerate(boxes):
        x = Inches(0.35 + i * 4.3)
        add_rect(s, x, Inches(6.25), Inches(4.0), Inches(1.08), MID_BG)
        add_rect(s, x, Inches(6.25), Inches(4.0), Inches(0.08), col)
        add_textbox(s, num, x + Inches(0.15), Inches(6.37), Inches(3.7), Inches(0.52),
                    font_size=30, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_textbox(s, label, x + Inches(0.15), Inches(6.87), Inches(3.7), Inches(0.38),
                    font_size=13, color=SUBTEXT, align=PP_ALIGN.CENTER)
    return s


def slide_model_b_hint_results(prs):
    """New slide — Hint scorer results (split out from the cluttered slide 18)."""
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    section_header_bar(s, "06  Model B — Hint Scorer Results (Logistic Regression)")

    add_textbox(s, "Evaluated on full test set (8,786 questions). Each question scored all article sentences; top-3 selected as hints.",
                Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.38), font_size=14, color=SUBTEXT)

    # Left: 4 big stat tiles
    add_textbox(s, "Classification Metrics", Inches(0.4), Inches(1.75), Inches(6.1), Inches(0.4),
                font_size=15, bold=True, color=ORANGE)

    tiles = [
        ("0.8681", "Accuracy",  LIGHT_GREY),
        ("0.2871", "Precision", ORANGE),
        ("0.8406", "Recall",    GREEN),
        ("0.4280", "F1",        LIGHT_GREY),
    ]
    for i, (val, name, col) in enumerate(tiles):
        col_i = i % 2; row_i = i // 2
        x = Inches(0.4 + col_i * 3.1)
        y = Inches(2.25 + row_i * 1.7)
        add_rect(s, x, y, Inches(2.85), Inches(1.5), MID_BG)
        add_rect(s, x, y, Inches(2.85), Inches(0.09), col)
        add_textbox(s, val, x + Inches(0.1), y + Inches(0.2), Inches(2.65), Inches(0.75),
                    font_size=36, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_textbox(s, name, x + Inches(0.1), y + Inches(0.98), Inches(2.65), Inches(0.38),
                    font_size=14, color=SUBTEXT, align=PP_ALIGN.CENTER)

    # Recall explanation box
    add_rect(s, Inches(0.4), Inches(5.7), Inches(6.1), Inches(1.55), RGBColor(0x0A, 0x3A, 0x18))
    add_rect(s, Inches(0.4), Inches(5.7), Inches(0.12), Inches(1.55), GREEN)
    add_textbox(s, "Why high recall is correct here",
                Inches(0.65), Inches(5.78), Inches(5.7), Inches(0.38),
                font_size=14, bold=True, color=GREEN)
    add_textbox(s,
        "It is better to surface a relevant sentence as a hint than to miss it.\n"
        "Low precision (0.29) is acceptable — users see 3 hints, not all candidates.\n"
        "False positives are less harmful than false negatives in hint generation.",
        Inches(0.65), Inches(6.18), Inches(5.7), Inches(0.95),
        font_size=13, color=LIGHT_GREY)

    add_rect(s, Inches(6.75), Inches(1.65), Inches(0.04), Inches(5.7), SUBTEXT)

    # Right: generation quality metrics
    add_textbox(s, "Generation Quality (Hint Text vs Reference)",
                Inches(7.0), Inches(1.75), Inches(6.1), Inches(0.4),
                font_size=15, bold=True, color=ORANGE)

    gen_metrics = [
        ("BLEU-1",  "0.2955", "Moderate exact n-gram match"),
        ("BLEU-4",  "0.2570", "4-gram precision"),
        ("ROUGE-1", "0.4380", "Unigram F1"),
        ("ROUGE-2", "0.3987", "Bigram F1"),
        ("ROUGE-L", "0.4337", "Longest common subsequence"),
        ("METEOR",  "0.6869", "Synonym + stem aware — highest score"),
    ]
    for i, (name, val, note) in enumerate(gen_metrics):
        y = Inches(2.25 + i * 0.73)
        add_rect(s, Inches(7.0), y, Inches(6.1), Inches(0.65), MID_BG)
        add_textbox(s, name, Inches(7.1), y + Inches(0.1), Inches(1.3), Inches(0.45),
                    font_size=14, bold=True, color=HIGHLIGHT)
        add_textbox(s, val, Inches(8.5), y + Inches(0.1), Inches(1.3), Inches(0.45),
                    font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, note, Inches(9.9), y + Inches(0.1), Inches(3.1), Inches(0.45),
                    font_size=12, italic=True, color=SUBTEXT)

    # Key insight
    add_rect(s, Inches(7.0), Inches(6.6), Inches(6.1), Inches(0.78), MID_BG)
    add_rect(s, Inches(7.0), Inches(6.6), Inches(0.12), Inches(0.78), HIGHLIGHT)
    add_textbox(s,
        "High METEOR (0.687) vs lower BLEU (0.296) confirms hints are semantically\n"
        "aligned with reference sentences but use different wording — appropriate\n"
        "since good hints paraphrase rather than copy.",
        Inches(7.2), Inches(6.65), Inches(5.8), Inches(0.68),
        font_size=13, color=LIGHT_GREY)
    return s


def slide_model_b_honest(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    section_header_bar(s, "06  Model B — Honest Evaluation: ModelBfinal")

    add_textbox(s,
        "ModelBfinal removes gold distractors from the candidate pool — forcing truly NOVEL generation:",
        Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.4), font_size=15, color=LIGHT_GREY)

    # comparison table
    rows = [
        ("Metric",   "Model_B (gold in pool)", "ModelBfinal (gen-only)", "Drop"),
        ("BLEU-1",   "0.8020",                 "0.0147",                 "−98.2%"),
        ("BLEU-4",   "0.6713",                 "0.0043",                 "−99.4%"),
        ("ROUGE-L",  "0.8060",                 "0.0450",                 "−94.4%"),
        ("METEOR",   "0.7618",                 "0.0238",                 "−96.9%"),
    ]
    for ri, row in enumerate(rows):
        y = Inches(1.9 + ri * 0.75)
        bg = ACCENT if ri == 0 else (MID_BG if ri % 2 else RGBColor(0x1E, 0x2D, 0x50))
        add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.68), bg)
        xs = [Inches(0.6), Inches(3.6), Inches(7.5), Inches(11.0)]
        ws = [Inches(3.0), Inches(3.9), Inches(3.5), Inches(1.8)]
        for ci, (cell, x, w) in enumerate(zip(row, xs, ws)):
            col = RED if ri > 0 and ci == 3 else (WHITE if ri == 0 else LIGHT_GREY)
            add_textbox(s, cell, x, y + Inches(0.1), w, Inches(0.5),
                        font_size=14, bold=(ri == 0), color=col)

    # Insight box
    add_rect(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.1), RGBColor(0x3A, 0x10, 0x10))
    add_rect(s, Inches(0.5), Inches(5.0), Inches(0.15), Inches(2.1), RED)
    add_textbox(s, "⚠  Critical Insight", Inches(0.8), Inches(5.1), Inches(11.8), Inches(0.45),
                font_size=17, bold=True, color=RED)
    add_textbox(s,
        "The ~97% drop in METEOR (0.762 → 0.024) reveals that Model B's impressive scores in the initial\n"
        "evaluation were largely inflated by re-ranking gold incorrect options that were already in the candidate pool.\n\n"
        "When forced to generate truly novel distractors, the model struggles — it cannot produce new text that\n"
        "maintains plausibility while being distinct from the correct answer. This is an honest and important\n"
        "limitation of the template-based candidate generation approach used in this project.",
        Inches(0.8), Inches(5.6), Inches(11.8), Inches(1.4), font_size=14, color=LIGHT_GREY)
    return s


def slide_results_summary(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    section_header_bar(s, "07  All Results — Consolidated Test-Set Metrics")

    add_textbox(s, "All metrics measured on RACE test set (8,786 questions):",
                Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.35), font_size=13, color=SUBTEXT)

    rows = [
        ("Component",                        "Algorithm",         "BLEU-1", "BLEU-4", "ROUGE-L", "METEOR", "Other"),
        ("Model A — Verifier",               "LinearSVC+Combined","0.4806", "0.3325", "0.4971",  "0.4400", "Q-Acc: 0.373"),
        ("Model A — QGen (baseline)",        "Template only",     "0.1185", "0.0238", "0.1700",  "0.0914", "—"),
        ("Model A — QGen v1 (ranker)",       "LinearSVC",         "0.1811", "0.0607", "0.2091",  "0.1598", "+53% METEOR↑"),
        ("Model A — QGen v2 (more templates)","LinearSVC v2",     "0.1758", "0.0597", "0.2055",  "0.1564", "Slight regress"),
        ("Model B — Distractor (with gold)", "Random Forest",     "0.8020", "0.6713", "0.8060",  "0.7618", "Inflated ⚠"),
        ("Model B — Distractor (gen-only)",  "Random Forest",     "0.0147", "0.0043", "0.0450",  "0.0238", "True gen"),
        ("Model B — Hints",                  "Logistic Reg.",     "0.2955", "0.2570", "0.4337",  "0.6869", "Recall: 0.84"),
    ]
    col_xs = [Inches(0.3), Inches(4.2), Inches(6.55), Inches(7.85), Inches(9.15), Inches(10.45), Inches(11.65)]
    col_ws = [Inches(3.85), Inches(2.3), Inches(1.25), Inches(1.25), Inches(1.25), Inches(1.18), Inches(1.5)]
    for ri, row in enumerate(rows):
        y = Inches(1.7 + ri * 0.66)
        special = ri in [4, 6]
        bg = ACCENT if ri == 0 else (RGBColor(0x06, 0x4A, 0x2C) if special else (MID_BG if ri % 2 else RGBColor(0x1E, 0x2D, 0x50)))
        add_rect(s, Inches(0.3), y, Inches(12.9), Inches(0.61), bg)
        for cell, x, w in zip(row, col_xs, col_ws):
            col = WHITE if ri == 0 else (GREEN if special else LIGHT_GREY)
            if ri == 5 and cell in ["0.8020", "0.6713", "0.8060", "0.7618"]:
                col = ORANGE
            add_textbox(s, cell, x + Inches(0.04), y + Inches(0.09), w, Inches(0.44),
                        font_size=12, bold=(ri == 0), color=col)

    add_textbox(s,
        "Green rows = honest best results.  Orange = inflated (gold in pool).  "
        "METEOR chosen as primary metric — it accounts for synonyms and stemming (more robust than BLEU).",
        Inches(0.3), Inches(7.1), Inches(12.9), Inches(0.35),
        font_size=12, italic=True, color=SUBTEXT)
    return s


def slide_metric_interpretation(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    section_header_bar(s, "07  Metric Interpretation — Why These Numbers?")

    points = [
        (HIGHLIGHT, "BLEU-4",
         "0.0607 for QGen (linear ranker)",
         "BLEU-4 requires 4 consecutive words to match. Generated questions use different phrasing "
         "from gold questions even when semantically correct. BLEU-4 of 0.06 is typical for "
         "template-based QGen vs human-written questions."),
        (GREEN, "METEOR",
         "0.160 for QGen / 0.687 for hints",
         "METEOR aligns synonyms and stemmed forms. Higher METEOR than BLEU confirms that "
         "generated text is semantically close to reference even without exact n-gram matches. "
         "QGen METEOR of 0.16 is in the expected range for non-neural template systems."),
        (ORANGE, "ROUGE-L",
         "0.209 for QGen / 0.434 for hints",
         "ROUGE-L measures longest common subsequence F1. The moderate ROUGE-L for hints (0.43) "
         "confirms selected sentences share significant subsequences with reference hints — "
         "appropriate since hints are drawn directly from the article."),
        (RED, "Q-level Accuracy",
         "0.373 for verifier (chance = 0.25)",
         "The verifier is correct 37.3% of the time on 4-way MCQ. A language-model approach "
         "would score ~0.70+. This is expected for a bag-of-words model — it captures surface "
         "overlap but not deep inference. Still provides a meaningful signal above random."),
        (RGBColor(0xAE, 0x81, 0xFF), "Silhouette",
         "0.006 (KMeans) vs 0.074 (GMM)",
         "Silhouette scores near zero indicate overlapping clusters in TF-IDF space — topics "
         "blend continuously. GMM on 11 numeric features shows cleaner separation (0.074) "
         "because cosine/overlap features have more interpretable Gaussian structure."),
    ]

    for i, (col, name, val, interp) in enumerate(points):
        y = Inches(1.35 + i * 1.18)
        add_rect(s, Inches(0.3), y, Inches(12.7), Inches(1.08), MID_BG)
        add_rect(s, Inches(0.3), y, Inches(0.14), Inches(1.08), col)
        add_textbox(s, name, Inches(0.55), y + Inches(0.06), Inches(2.0), Inches(0.38),
                    font_size=16, bold=True, color=col)
        add_textbox(s, val, Inches(2.65), y + Inches(0.06), Inches(4.5), Inches(0.38),
                    font_size=14, bold=True, color=WHITE)
        add_textbox(s, interp, Inches(0.55), y + Inches(0.48), Inches(12.0), Inches(0.55),
                    font_size=13, color=LIGHT_GREY)
    return s


def slide_system_architecture(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    section_header_bar(s, "08  System Architecture & Inference Flow")

    # Two paths
    add_textbox(s, "Path 1 — RACE article (loaded from dataset)",
                Inches(0.4), Inches(1.3), Inches(5.8), Inches(0.4),
                font_size=14, bold=True, color=HIGHLIGHT)
    path1 = [
        "User loads article from RACE tab",
        "lookup_sample_by_article() →\nfinds row in test.csv",
        "Gold Q&A returned directly\n(no QGen inference needed)",
        "Model B generates\nfresh distractors",
        "Model B generates\n3 graduated hints",
    ]
    for i, step in enumerate(path1):
        x = Inches(0.35 + i * 2.55)
        add_rect(s, x, Inches(1.82), Inches(2.35), Inches(1.1), ACCENT)
        add_textbox(s, step, x + Inches(0.08), Inches(1.87), Inches(2.2), Inches(1.0),
                    font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
        if i < 4:
            add_textbox(s, "→", x + Inches(2.35), Inches(2.17), Inches(0.25), Inches(0.38),
                        font_size=16, bold=True, color=HIGHLIGHT, align=PP_ALIGN.CENTER)

    add_rect(s, 0, Inches(3.1), SLIDE_W, Inches(0.04), SUBTEXT)

    add_textbox(s, "Path 2 — Custom article (user pasted)",
                Inches(0.4), Inches(3.25), Inches(5.8), Inches(0.4),
                font_size=14, bold=True, color=GREEN)
    path2 = [
        "User pastes custom\narticle text",
        "lookup returns None →\nModel A QGen pipeline",
        "Top-2 sentences ×\ntemplates → 27 features",
        "LinearSVC ranker\npicks best question",
        "Model B generates\ndistractors + hints",
    ]
    for i, step in enumerate(path2):
        x = Inches(0.35 + i * 2.55)
        add_rect(s, x, Inches(3.72), Inches(2.35), Inches(1.1), RGBColor(0x0A, 0x3A, 0x18))
        add_textbox(s, step, x + Inches(0.08), Inches(3.77), Inches(2.2), Inches(1.0),
                    font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
        if i < 4:
            add_textbox(s, "→", x + Inches(2.35), Inches(4.07), Inches(0.25), Inches(0.38),
                        font_size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    add_rect(s, 0, Inches(5.0), SLIDE_W, Inches(0.04), SUBTEXT)

    # Always happens
    add_textbox(s, "Always (both paths) — when user checks answer:",
                Inches(0.4), Inches(5.12), Inches(12.5), Inches(0.38),
                font_size=14, bold=True, color=ORANGE)
    always = [
        ("Model A Verifier", "LinearSVC predicts if\nuser's choice is correct"),
        ("BLEU/ROUGE/METEOR", "Computed vs gold answer\n→ logged to dashboard"),
        ("Session Log", "Entry added:\nq, answer, metrics, latency"),
    ]
    for i, (title, desc) in enumerate(always):
        x = Inches(0.4 + i * 4.3)
        add_rect(s, x, Inches(5.58), Inches(4.1), Inches(1.0), MID_BG)
        add_rect(s, x, Inches(5.58), Inches(4.1), Inches(0.06), ORANGE)
        add_textbox(s, title, x + Inches(0.1), Inches(5.68), Inches(3.9), Inches(0.38),
                    font_size=14, bold=True, color=ORANGE)
        add_textbox(s, desc, x + Inches(0.1), Inches(6.07), Inches(3.9), Inches(0.45),
                    font_size=13, color=LIGHT_GREY)
    return s


def slide_ui_screens(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    accent_bar(s)
    section_header_bar(s, "08  Streamlit UI — 4 Screens")

    screens = [
        ("📄  Article Input",
         "• Paste custom article OR load from RACE dataset tab\n"
         "• Minimum 30 characters validation\n"
         "• 'Generate Quiz' triggers full Model A+B inference\n"
         "• Live character + word counter",
         HIGHLIGHT),
        ("❓  Quiz View",
         "• Displays generated question + 4 options (A/B/C/D)\n"
         "• User selects answer → Model A verifier runs\n"
         "• Correct/Incorrect feedback with confidence score\n"
         "• BLEU/ROUGE/METEOR computed + logged",
         GREEN),
        ("💡  Hint Panel",
         "• 3 graduated hints revealed one at a time\n"
         "• 'Reveal Answer' locked until all 3 hints shown\n"
         "• Re-read passage expander always available\n"
         "• Reset hints button",
         ORANGE),
        ("📊  Dashboard",
         "• Session BLEU/ROUGE/METEOR averages\n"
         "• Latency tracking (avg/min/max, 10s SLA check)\n"
         "• Session log table with per-question metrics\n"
         "• CSV export + clear log",
         RED),
    ]

    for i, (title, body, col) in enumerate(screens):
        col_i = i % 2
        row_i = i // 2
        x = Inches(0.3 + col_i * 6.55)
        y = Inches(1.35 + row_i * 2.85)
        add_rect(s, x, y, Inches(6.3), Inches(2.65), MID_BG)
        add_rect(s, x, y, Inches(6.3), Inches(0.52), col)
        add_textbox(s, title, x + Inches(0.15), y + Inches(0.08), Inches(6.0), Inches(0.4),
                    font_size=17, bold=True, color=DARK_BG)
        add_textbox(s, body, x + Inches(0.15), y + Inches(0.62), Inches(6.0), Inches(1.9),
                    font_size=14, color=LIGHT_GREY)
    return s


def slide_closing(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.12), HIGHLIGHT)
    add_rect(s, 0, Inches(7.38), SLIDE_W, Inches(0.12), HIGHLIGHT)

    add_textbox(s, "Key Takeaways", Inches(1.0), Inches(1.0), Inches(11.3), Inches(0.65),
                font_size=34, bold=True, color=HIGHLIGHT, align=PP_ALIGN.CENTER)

    takeaways = [
        (HIGHLIGHT, "Combined features beat TF-IDF alone",
         "11 lexical/cosine features + 50K TF-IDF → F1 0.517 vs 0.503 TF-IDF only"),
        (GREEN, "QGen ranker = +53% METEOR over baseline",
         "LinearSVC ranker on 27 features meaningfully improves template selection"),
        (ORANGE, "More templates ≠ better (v2 regressed)",
         "Diversity without a stronger signal dilutes ranking quality"),
        (RED, "Honest evaluation matters",
         "Distractor METEOR drops 97% when gold options removed — always evaluate with gold excluded"),
        (RGBColor(0xAE, 0x81, 0xFF), "Traditional ML has limits on this task",
         "37% verifier accuracy (vs 25% chance) shows surface overlap is not sufficient for deep comprehension"),
    ]

    for i, (col, title, body) in enumerate(takeaways):
        y = Inches(1.9 + i * 1.05)
        add_rect(s, Inches(0.8), y, Inches(11.7), Inches(0.9), MID_BG)
        add_rect(s, Inches(0.8), y, Inches(0.14), Inches(0.9), col)
        add_textbox(s, title, Inches(1.1), y + Inches(0.05), Inches(5.5), Inches(0.38),
                    font_size=15, bold=True, color=col)
        add_textbox(s, body, Inches(1.1), y + Inches(0.48), Inches(10.8), Inches(0.36),
                    font_size=13, color=LIGHT_GREY)

    add_textbox(s, "Thank you — Demo Q&A",
                Inches(1.0), Inches(7.0), Inches(11.3), Inches(0.38),
                font_size=18, bold=True, color=SUBTEXT, align=PP_ALIGN.CENTER)
    return s


# ══════════════════════════════════════════════════════════════════════════════
# BUILD
# ══════════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()

    slides_fns = [
        slide_title,
        slide_toc,
        slide_dataset_overview,
        slide_dataset_expansion,
        slide_eda_distributions,
        slide_eda_template_wins,
        slide_eda_cosine,
        slide_eda_tfidf_vs_onehot,
        slide_statistical_analysis,
        slide_stat_pseudo_labels,
        slide_preprocessing,
        slide_preprocessing_features,
        slide_model_a_selection,
        slide_model_a_training,
        slide_model_a_qgen_features,
        slide_model_a_results,
        slide_model_b_overview,
        slide_model_b_features,
        slide_model_b_training,
        slide_model_b_hint_results,
        slide_model_b_honest,
        slide_results_summary,
        slide_metric_interpretation,
        slide_system_architecture,
        slide_ui_screens,
        slide_closing,
    ]

    total = len(slides_fns)
    built = []
    for fn in slides_fns:
        sl = fn(prs)
        built.append(sl)

    # Add slide numbers (skip title)
    for i, sl in enumerate(built[1:], start=2):
        slide_number(sl, i, total, prs)

    out = "notebooks/RACE_Quiz_System_Presentation.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({total} slides)")


if __name__ == "__main__":
    build()
